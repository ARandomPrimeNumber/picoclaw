import os
import sys
import subprocess

def install_and_import(package, import_name=None):
    if import_name is None:
        import_name = package
    try:
        __import__(import_name)
    except ImportError:
        # Silently install missing packages
        subprocess.check_call([sys.executable, "-m", "pip", "install", package, "-q", "--disable-pip-version-check"])
        __import__(import_name)

# Ensure dependencies exist
install_and_import("python-docx", "docx")
install_and_import("openpyxl", "openpyxl")
install_and_import("python-pptx", "pptx")
install_and_import("PyPDF2", "PyPDF2")

import docx
import openpyxl
import pptx
import PyPDF2
import csv

def extract_docx(path):
    doc = docx.Document(path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    # Encode to utf-8 to handle special characters
    return "\n".join(paragraphs).encode('utf-8', errors='replace').decode('utf-8')

def extract_xlsx(path):
    wb = openpyxl.load_workbook(path, data_only=True)
    out = []
    out.append(f"Spreadsheet contains {len(wb.sheetnames)} sheets: {', '.join(wb.sheetnames)}")
    out.append("---")
    
    for sheetname in wb.sheetnames:
        sheet = wb[sheetname]
        out.append(f"\n[ SHEET: {sheetname} ]")
        
        # Cap at 100 rows to avoid context overflow
        rows_processed = 0
        for row in sheet.iter_rows(values_only=True):
            if rows_processed >= 100:
                out.append(f"...[truncated, {sheet.max_row - 100} more rows]")
                break
                
            row_data = [str(cell) if cell is not None else "" for cell in row]
            if any(row_data): # Skip completely empty rows
                out.append(" | ".join(row_data))
                rows_processed += 1
                
    return "\n".join(out)

def extract_pptx(path):
    prs = pptx.Presentation(path)
    out = []
    out.append(f"Presentation contains {len(prs.slides)} slides.")
    out.append("---")
    
    for i, slide in enumerate(prs.slides):
        out.append(f"\n[ SLIDE {i+1} ]")
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        if slide_text:
            out.append("\n".join(slide_text))
        else:
            out.append("(No text on slide)")
            
    return "\n".join(out)

def extract_pdf(path):
    out = []
    with open(path, 'rb') as f:
        reader = PyPDF2.PdfReader(f)
        out.append(f"PDF contains {len(reader.pages)} pages.")
        out.append("---")
        
        # Read up to 20 pages max to avoid massive prompt overflow
        max_pages = min(20, len(reader.pages))
        for i in range(max_pages):
            out.append(f"\n[ PAGE {i+1} ]")
            text = reader.pages[i].extract_text()
            if text:
                out.append(text.strip())
                
        if len(reader.pages) > 20:
            out.append(f"\n...[truncated, {len(reader.pages) - 20} more pages]")
            
    return "\n".join(out)

def extract_csv(path):
    out = []
    # Try different encodings
    encodings = ['utf-8', 'utf-8-sig', 'latin1', 'cp1252']
    
    for enc in encodings:
        try:
            with open(path, 'r', encoding=enc) as f:
                reader = csv.reader(f)
                rows_processed = 0
                for row in reader:
                    if rows_processed >= 200:
                        out.append("...[truncated remaining rows]")
                        break
                    if any(row):
                        out.append(" | ".join(row))
                        rows_processed += 1
            return "\n".join(out)
        except UnicodeDecodeError:
            continue
            
    return "Error: Could not decode CSV file."

def main():
    import io
    # Set stdout to UTF-8 to handle special characters
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
    
    if len(sys.argv) < 2:
        print("Usage: python extract_text.py <filepath>", file=sys.stderr)
        sys.exit(1)
        
    path = sys.argv[1]
    
    # Check if inside workspace (sandbox)
    abs_path = os.path.abspath(path)
    if not os.path.exists(abs_path):
        print(f"Error: File not found: {abs_path}", file=sys.stderr)
        sys.exit(1)
        
    ext = os.path.splitext(abs_path)[1].lower()
    
    try:
        if ext == '.docx':
            print(extract_docx(abs_path))
        elif ext == '.xlsx':
            print(extract_xlsx(abs_path))
        elif ext == '.pptx':
            print(extract_pptx(abs_path))
        elif ext == '.pdf':
            print(extract_pdf(abs_path))
        elif ext == '.csv':
            print(extract_csv(abs_path))
        elif ext in ['.txt', '.md', '.log', '.json']:
            # For plain text files, standard read file is better, but handle it anyway if passed
            with open(abs_path, 'r', encoding='utf-8', errors='ignore') as f:
                # Truncate at 1000 lines
                lines = [next(f) for _ in range(1000)]
                print("".join(lines))
        else:
            print(f"Error: Unsupported file extension '{ext}'", file=sys.stderr)
            sys.exit(1)
    except Exception as e:
        print(f"Error extracting {ext} file: {str(e)}", file=sys.stderr)
        sys.exit(1)

if __name__ == "__main__":
    main()
